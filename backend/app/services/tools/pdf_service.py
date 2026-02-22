"""PDF processing service."""

import io
import os
import tempfile
from concurrent.futures import ThreadPoolExecutor
from typing import Optional, Tuple

from pypdf import PdfReader, PdfWriter

from app.core.exceptions import BadRequestError, FileProcessingError

# Use CPU count for better parallelization (min 4, max 16)
_max_workers = min(max(os.cpu_count() or 4, 4), 16)
_executor = ThreadPoolExecutor(max_workers=_max_workers)


class PDFService:
    """Service for PDF processing operations."""

    async def split(
        self,
        pdf_bytes: bytes,
        page_ranges: str,
    ) -> Tuple[list[Tuple[bytes, list[int]]], int]:
        """Split PDF by page ranges."""
        try:
            reader = PdfReader(io.BytesIO(pdf_bytes))
            total_pages = len(reader.pages)

            # Parse page ranges
            ranges = self._parse_page_ranges(page_ranges, total_pages)

            result_files = []

            for page_range in ranges:
                writer = PdfWriter()

                for page_num in page_range:
                    writer.add_page(reader.pages[page_num - 1])  # 1-indexed to 0-indexed

                # Compress output
                writer.compress_identical_objects(remove_identicals=True, remove_orphans=True)

                output = io.BytesIO()
                writer.write(output)
                result_files.append((output.getvalue(), page_range))

            return result_files, total_pages

        except Exception as e:
            raise FileProcessingError(message=f"PDF split failed: {str(e)}")

    async def merge(self, pdf_contents: list[bytes]) -> Tuple[bytes, int]:
        """Merge multiple PDFs into one."""
        try:
            writer = PdfWriter()
            total_pages = 0

            for content in pdf_contents:
                reader = PdfReader(io.BytesIO(content))
                for page in reader.pages:
                    writer.add_page(page)
                    total_pages += 1

            # Compress output
            writer.compress_identical_objects(remove_identicals=True, remove_orphans=True)

            output = io.BytesIO()
            writer.write(output)

            return output.getvalue(), total_pages

        except Exception as e:
            raise FileProcessingError(message=f"PDF merge failed: {str(e)}")

    async def compress(
        self,
        pdf_bytes: bytes,
        compression_level: str,
    ) -> Tuple[bytes, int]:
        """Compress PDF to reduce file size."""
        try:
            import asyncio

            # Use Ghostscript for real compression if available
            loop = asyncio.get_event_loop()
            result = await loop.run_in_executor(
                _executor,
                lambda: self._compress_with_ghostscript(pdf_bytes, compression_level)
            )

            if result is not None:
                compressed_bytes, page_count = result
                # Only return compressed if it's actually smaller
                if len(compressed_bytes) < len(pdf_bytes):
                    return compressed_bytes, page_count

            # Fallback: try pypdf compression
            reader = PdfReader(io.BytesIO(pdf_bytes))
            writer = PdfWriter()

            for page in reader.pages:
                page.compress_content_streams()
                writer.add_page(page)

            # Remove duplication and compress
            writer.compress_identical_objects(remove_identicals=True, remove_orphans=True)

            output = io.BytesIO()
            writer.write(output)
            compressed = output.getvalue()

            # Only return compressed version if it's actually smaller
            if len(compressed) < len(pdf_bytes):
                return compressed, len(reader.pages)
            else:
                # Return original if compression didn't help
                return pdf_bytes, len(reader.pages)

        except Exception as e:
            raise FileProcessingError(message=f"PDF compression failed: {str(e)}")

    def _compress_with_ghostscript(
        self, pdf_bytes: bytes, compression_level: str
    ) -> Tuple[bytes, int] | None:
        """Compress PDF using Ghostscript for better results."""
        import subprocess
        import shutil

        # Check if Ghostscript is available
        if not shutil.which("gs"):
            return None

        # Map compression levels to Ghostscript settings
        gs_settings = {
            "low": "/prepress",      # High quality, larger file
            "medium": "/ebook",       # Medium quality, good balance
            "high": "/screen",        # Lower quality, smallest file
        }

        pdfsetting = gs_settings.get(compression_level, "/ebook")

        try:
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as input_file:
                input_file.write(pdf_bytes)
                input_path = input_file.name

            output_path = input_path.replace(".pdf", "_compressed.pdf")

            try:
                # Run Ghostscript
                result = subprocess.run(
                    [
                        "gs",
                        "-sDEVICE=pdfwrite",
                        "-dCompatibilityLevel=1.4",
                        f"-dPDFSETTINGS={pdfsetting}",
                        "-dNOPAUSE",
                        "-dQUIET",
                        "-dBATCH",
                        "-dCompressFonts=true",
                        "-dSubsetFonts=true",
                        "-dColorImageDownsampleType=/Bicubic",
                        "-dGrayImageDownsampleType=/Bicubic",
                        "-dMonoImageDownsampleType=/Bicubic",
                        f"-sOutputFile={output_path}",
                        input_path,
                    ],
                    capture_output=True,
                    timeout=60,
                )

                if result.returncode == 0 and os.path.exists(output_path):
                    with open(output_path, "rb") as f:
                        compressed_bytes = f.read()

                    # Get page count
                    reader = PdfReader(io.BytesIO(compressed_bytes))
                    return compressed_bytes, len(reader.pages)

                return None

            finally:
                if os.path.exists(input_path):
                    os.unlink(input_path)
                if os.path.exists(output_path):
                    os.unlink(output_path)

        except Exception:
            return None

    async def to_word(self, pdf_bytes: bytes) -> Tuple[bytes, int]:
        """Convert PDF to Word document."""
        import asyncio

        # pdf2docx requires file paths, so use temp files
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as pdf_file:
            pdf_file.write(pdf_bytes)
            pdf_path = pdf_file.name

        docx_path = pdf_path.replace(".pdf", ".docx")

        try:
            # Run in thread pool (CPU-bound)
            loop = asyncio.get_event_loop()
            success = await loop.run_in_executor(
                _executor,
                lambda: self._convert_pdf_to_docx(pdf_path, docx_path)
            )

            if not success or not os.path.exists(docx_path):
                raise FileProcessingError(
                    message="PDF conversion failed. The PDF may be corrupted or use unsupported features."
                )

            # Read result
            with open(docx_path, "rb") as f:
                docx_bytes = f.read()

            # Get page count using fitz (more reliable)
            import fitz
            doc = fitz.open(pdf_path)
            total_pages = len(doc)
            doc.close()

            return docx_bytes, total_pages

        except FileProcessingError:
            raise
        except Exception as e:
            raise FileProcessingError(message=f"PDF to Word conversion failed: {str(e)}")
        finally:
            # Cleanup temp files
            if os.path.exists(pdf_path):
                os.unlink(pdf_path)
            if os.path.exists(docx_path):
                os.unlink(docx_path)

    def _convert_pdf_to_docx(self, pdf_path: str, docx_path: str) -> bool:
        """Synchronous PDF to DOCX conversion with fallback."""
        # Try pdf2docx first
        try:
            from pdf2docx import Converter
            cv = Converter(pdf_path)
            cv.convert(docx_path)
            cv.close()
            return True
        except Exception:
            pass

        # Fallback: Use pymupdf to extract text and create a simple docx
        try:
            import fitz
            from docx import Document
            from docx.shared import Pt, Inches

            doc = fitz.open(pdf_path)
            word_doc = Document()

            for page_num in range(len(doc)):
                page = doc[page_num]

                # Add page break between pages (except first)
                if page_num > 0:
                    word_doc.add_page_break()

                # Extract text blocks
                blocks = page.get_text("blocks")

                for block in blocks:
                    if block[6] == 0:  # Text block
                        text = block[4].strip()
                        if text:
                            para = word_doc.add_paragraph(text)
                            para.style.font.size = Pt(11)

                # Extract images
                images = page.get_images()
                for img_index, img in enumerate(images):
                    try:
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]

                        # Save image temporarily
                        img_path = f"{pdf_path}_img_{page_num}_{img_index}.{base_image['ext']}"
                        with open(img_path, "wb") as img_file:
                            img_file.write(image_bytes)

                        # Add to document
                        word_doc.add_picture(img_path, width=Inches(5))

                        # Cleanup
                        os.unlink(img_path)
                    except Exception:
                        continue

            doc.close()
            word_doc.save(docx_path)
            return True

        except Exception:
            return False

    def _parse_page_ranges(self, ranges_str: str, total_pages: int) -> list[list[int]]:
        """Parse page range string like '1-5,6-10' or 'all'."""
        if ranges_str.lower() == "all":
            # Return each page as separate range
            return [[i] for i in range(1, total_pages + 1)]

        ranges = []

        for part in ranges_str.split(","):
            part = part.strip()

            if "-" in part:
                try:
                    start, end = part.split("-")
                    start = int(start.strip())
                    end = int(end.strip())

                    if start < 1 or end > total_pages or start > end:
                        raise BadRequestError(
                            message=f"Invalid page range: {part}. PDF has {total_pages} pages."
                        )

                    ranges.append(list(range(start, end + 1)))
                except ValueError:
                    raise BadRequestError(message=f"Invalid page range format: {part}")
            else:
                try:
                    page = int(part)
                    if page < 1 or page > total_pages:
                        raise BadRequestError(
                            message=f"Invalid page number: {page}. PDF has {total_pages} pages."
                        )
                    ranges.append([page])
                except ValueError:
                    raise BadRequestError(message=f"Invalid page number: {part}")

        if not ranges:
            raise BadRequestError(message="No valid page ranges specified")

        return ranges

    async def remove_watermark(
        self,
        pdf_bytes: bytes,
        watermark_text: str = "",
        remove_images: bool = False,
        remove_text_watermarks: bool = True,
    ) -> Tuple[bytes, int, dict]:
        """
        Attempt to remove watermarks from PDF.
        Returns (cleaned_pdf_bytes, page_count, removal_stats).

        This works best on:
        - Text watermarks added as overlays
        - Annotation-based watermarks
        - Image watermarks in separate layers (if remove_images=True)

        May NOT work on:
        - Watermarks baked into page images
        - Heavily embedded watermarks in content streams
        """
        import fitz

        try:
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            total_pages = len(doc)

            stats = {
                "annotations_removed": 0,
                "text_instances_removed": 0,
                "images_removed": 0,
            }

            # Common watermark patterns to look for
            watermark_patterns = [
                "watermark", "draft", "confidential", "sample",
                "copy", "void", "preview", "demo", "trial",
                "do not copy", "not for distribution"
            ]

            # Add user-specified watermark text
            if watermark_text:
                watermark_patterns.insert(0, watermark_text.lower())

            for page in doc:
                # 1. Remove watermark annotations
                annots_to_remove = []
                for annot in page.annots() or []:
                    annot_info = annot.info
                    content = annot_info.get("content", "").lower()
                    title = annot_info.get("title", "").lower()

                    # Check if annotation looks like a watermark
                    is_watermark = any(
                        pattern in content or pattern in title
                        for pattern in watermark_patterns
                    )

                    # Also check for stamp/watermark annotation types
                    if annot.type[0] in [13, 14] or is_watermark:  # Stamp, Watermark types
                        annots_to_remove.append(annot)

                for annot in annots_to_remove:
                    page.delete_annot(annot)
                    stats["annotations_removed"] += 1

                # 2. Try to remove text watermarks from content
                if remove_text_watermarks:
                    text_dict = page.get_text("dict")

                    for block in text_dict.get("blocks", []):
                        if block.get("type") == 0:  # Text block
                            for line in block.get("lines", []):
                                for span in line.get("spans", []):
                                    text = span.get("text", "").lower()

                                    # Check if text matches watermark pattern
                                    if any(pattern in text for pattern in watermark_patterns):
                                        # Try to redact this text
                                        try:
                                            rect = fitz.Rect(span.get("bbox"))
                                            # Add white redaction to cover the watermark
                                            page.add_redact_annot(rect, fill=(1, 1, 1))
                                            stats["text_instances_removed"] += 1
                                        except Exception:
                                            pass

                    # Apply all redactions
                    page.apply_redactions()

                # 3. Optionally remove images that might be watermarks
                if remove_images:
                    images = page.get_images()
                    for img in images:
                        try:
                            xref = img[0]
                            # Get image info
                            img_info = doc.extract_image(xref)
                            img_bytes = img_info.get("image", b"")

                            # Heuristic: small semi-transparent images are often watermarks
                            # This is imperfect but can help
                            width = img_info.get("width", 0)
                            height = img_info.get("height", 0)

                            # Skip main content images (larger than 200x200)
                            if width > 200 and height > 200:
                                continue

                            # Try to remove the image
                            page.delete_image(xref)
                            stats["images_removed"] += 1
                        except Exception:
                            pass

            # Save the modified PDF with compression
            output = io.BytesIO()
            doc.save(output, garbage=4, deflate=True, clean=True)
            doc.close()

            return output.getvalue(), total_pages, stats

        except Exception as e:
            raise FileProcessingError(message=f"Watermark removal failed: {str(e)}")

    def apply_watermark(self, pdf_bytes: bytes) -> bytes:
        """Apply watermark to PDF (for free tier)."""
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.colors import Color

            # Create watermark PDF
            watermark_buffer = io.BytesIO()
            c = canvas.Canvas(watermark_buffer, pagesize=letter)

            # Set watermark style
            c.setFillColor(Color(0.7, 0.7, 0.7, alpha=0.3))
            c.setFont("Helvetica", 40)

            # Rotate and position
            c.saveState()
            c.translate(300, 400)
            c.rotate(45)
            c.drawCentredString(0, 0, "ToolHub Free")
            c.restoreState()

            c.save()

            # Merge watermark with original PDF
            watermark_buffer.seek(0)
            watermark_reader = PdfReader(watermark_buffer)
            watermark_page = watermark_reader.pages[0]

            reader = PdfReader(io.BytesIO(pdf_bytes))
            writer = PdfWriter()

            for page in reader.pages:
                page.merge_page(watermark_page)
                writer.add_page(page)

            output = io.BytesIO()
            writer.write(output)

            return output.getvalue()

        except Exception:
            # If watermarking fails, return original
            return pdf_bytes

    async def to_images(
        self,
        pdf_bytes: bytes,
        image_format: str = "jpeg",
        quality: int = 95,
        dpi: int = 200,
    ) -> Tuple[list[bytes], int]:
        """
        Convert PDF pages to images (JPG or PNG).

        Args:
            pdf_bytes: PDF file bytes
            image_format: "jpeg" or "png"
            quality: Image quality (1-100, for JPEG)
            dpi: Resolution in DPI

        Returns:
            Tuple of (list of image bytes, total pages)
        """
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            total_pages = len(doc)
            images = []

            # Calculate zoom factor from DPI (72 DPI is default)
            zoom = dpi / 72
            mat = fitz.Matrix(zoom, zoom)

            for page_num in range(total_pages):
                page = doc[page_num]

                # Render page to pixmap
                pix = page.get_pixmap(matrix=mat, alpha=False)

                # Convert to desired format
                if image_format.lower() == "png":
                    img_bytes = pix.tobytes("png")
                else:  # jpeg
                    img_bytes = pix.tobytes("jpeg", jpg_quality=quality)

                images.append(img_bytes)

            doc.close()
            return images, total_pages

        except Exception as e:
            raise FileProcessingError(message=f"PDF to image conversion failed: {str(e)}")

    async def from_images(
        self,
        image_contents: list[bytes],
        page_size: str = "A4",
        orientation: str = "portrait",
        margin: int = 0,
    ) -> Tuple[bytes, int]:
        """
        Convert images (JPG, PNG, etc.) to PDF.

        Args:
            image_contents: List of image file bytes
            page_size: "A4", "Letter", "Legal", or "auto" (fit to image)
            orientation: "portrait" or "landscape"
            margin: Margin in points (0-100)

        Returns:
            Tuple of (PDF bytes, number of pages)
        """
        try:
            from PIL import Image
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import A4, LETTER, LEGAL, landscape
            from reportlab.lib.utils import ImageReader

            # Determine page size
            page_sizes = {
                "a4": A4,
                "letter": LETTER,
                "legal": LEGAL,
            }

            output = io.BytesIO()
            c = None
            total_pages = len(image_contents)

            for idx, img_bytes in enumerate(image_contents):
                # Load image to get dimensions
                img = Image.open(io.BytesIO(img_bytes))
                img_width, img_height = img.size

                # Determine page size
                if page_size.lower() == "auto":
                    # Fit to image size
                    if orientation == "landscape":
                        pdf_width = max(img_width, img_height) + (margin * 2)
                        pdf_height = min(img_width, img_height) + (margin * 2)
                    else:
                        pdf_width = min(img_width, img_height) + (margin * 2)
                        pdf_height = max(img_width, img_height) + (margin * 2)
                    page_dims = (pdf_width, pdf_height)
                else:
                    page_dims = page_sizes.get(page_size.lower(), A4)
                    if orientation == "landscape":
                        page_dims = landscape(page_dims)

                # Create canvas for first page or new page for subsequent
                if c is None:
                    c = canvas.Canvas(output, pagesize=page_dims)
                else:
                    c.setPageSize(page_dims)
                    c.showPage()

                # Calculate image position to center it with margins
                page_width, page_height = page_dims
                available_width = page_width - (margin * 2)
                available_height = page_height - (margin * 2)

                # Scale image to fit while maintaining aspect ratio
                img_aspect = img_width / img_height
                page_aspect = available_width / available_height

                if img_aspect > page_aspect:
                    # Image is wider - fit to width
                    draw_width = available_width
                    draw_height = available_width / img_aspect
                else:
                    # Image is taller - fit to height
                    draw_height = available_height
                    draw_width = available_height * img_aspect

                # Center image on page
                x = (page_width - draw_width) / 2
                y = (page_height - draw_height) / 2

                # Draw image
                img_reader = ImageReader(io.BytesIO(img_bytes))
                c.drawImage(img_reader, x, y, width=draw_width, height=draw_height)

            if c:
                c.save()

            return output.getvalue(), total_pages

        except Exception as e:
            raise FileProcessingError(message=f"Image to PDF conversion failed: {str(e)}")

    async def rotate(
        self,
        pdf_bytes: bytes,
        rotation: int,
        pages: str = "all",
    ) -> Tuple[bytes, int]:
        """
        Rotate PDF pages.

        Args:
            pdf_bytes: PDF file bytes
            rotation: Rotation angle (90, 180, 270, or -90)
            pages: Page selection - "all", "even", "odd", or page numbers like "1,3,5-7"

        Returns:
            Tuple of (rotated PDF bytes, total pages)
        """
        try:
            reader = PdfReader(io.BytesIO(pdf_bytes))
            writer = PdfWriter()
            total_pages = len(reader.pages)

            # Parse page selection
            if pages.lower() == "all":
                pages_to_rotate = set(range(total_pages))
            elif pages.lower() == "even":
                pages_to_rotate = set(range(1, total_pages, 2))  # 0-indexed: 1, 3, 5...
            elif pages.lower() == "odd":
                pages_to_rotate = set(range(0, total_pages, 2))  # 0-indexed: 0, 2, 4...
            else:
                # Parse specific pages like "1,3,5-7"
                pages_to_rotate = self._parse_page_selection(pages, total_pages)

            # Rotate pages
            for idx, page in enumerate(reader.pages):
                if idx in pages_to_rotate:
                    page.rotate(rotation)
                writer.add_page(page)

            # Compress output
            writer.compress_identical_objects(remove_identicals=True, remove_orphans=True)

            output = io.BytesIO()
            writer.write(output)

            return output.getvalue(), total_pages

        except Exception as e:
            raise FileProcessingError(message=f"PDF rotation failed: {str(e)}")

    async def unlock(self, pdf_bytes: bytes, password: str) -> Tuple[bytes, int, bool]:
        """
        Remove password protection from PDF.

        Args:
            pdf_bytes: Encrypted PDF file bytes
            password: Password to unlock the PDF

        Returns:
            Tuple of (unlocked PDF bytes, total pages, was_encrypted)
        """
        try:
            reader = PdfReader(io.BytesIO(pdf_bytes))
            was_encrypted = reader.is_encrypted

            if was_encrypted:
                # Try to decrypt with password
                if not reader.decrypt(password):
                    raise BadRequestError(message="Invalid password")

            # Create new PDF without encryption
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)

            # Compress output
            writer.compress_identical_objects(remove_identicals=True, remove_orphans=True)

            output = io.BytesIO()
            writer.write(output)

            return output.getvalue(), len(reader.pages), was_encrypted

        except BadRequestError:
            raise
        except Exception as e:
            raise FileProcessingError(message=f"PDF unlock failed: {str(e)}")

    async def protect(
        self,
        pdf_bytes: bytes,
        user_password: str,
        owner_password: Optional[str] = None,
        allow_printing: bool = True,
        allow_modification: bool = False,
        allow_copying: bool = False,
    ) -> Tuple[bytes, int]:
        """
        Add password protection to PDF.

        Args:
            pdf_bytes: PDF file bytes
            user_password: Password required to open the PDF
            owner_password: Password for full permissions (defaults to user_password)
            allow_printing: Allow printing
            allow_modification: Allow content modification
            allow_copying: Allow text/image copying

        Returns:
            Tuple of (protected PDF bytes, total pages)
        """
        try:
            reader = PdfReader(io.BytesIO(pdf_bytes))
            writer = PdfWriter()

            for page in reader.pages:
                writer.add_page(page)

            # Apply encryption
            writer.encrypt(
                user_password=user_password,
                owner_password=owner_password or user_password,
                permissions_flag=self._build_permissions_flag(
                    allow_printing, allow_modification, allow_copying
                ),
            )

            # Compress output
            writer.compress_identical_objects(remove_identicals=True, remove_orphans=True)

            output = io.BytesIO()
            writer.write(output)

            return output.getvalue(), len(reader.pages)

        except Exception as e:
            raise FileProcessingError(message=f"PDF protection failed: {str(e)}")

    def _parse_page_selection(self, selection: str, total_pages: int) -> set[int]:
        """Parse page selection string like '1,3,5-7' into 0-indexed page numbers."""
        pages = set()

        for part in selection.split(","):
            part = part.strip()

            if "-" in part:
                start, end = part.split("-")
                start = int(start.strip()) - 1  # Convert to 0-indexed
                end = int(end.strip()) - 1

                if start < 0 or end >= total_pages or start > end:
                    raise BadRequestError(
                        message=f"Invalid page range: {part}. PDF has {total_pages} pages."
                    )

                pages.update(range(start, end + 1))
            else:
                page = int(part) - 1  # Convert to 0-indexed
                if page < 0 or page >= total_pages:
                    raise BadRequestError(
                        message=f"Invalid page number: {part}. PDF has {total_pages} pages."
                    )
                pages.add(page)

        return pages

    def _build_permissions_flag(
        self, allow_printing: bool, allow_modification: bool, allow_copying: bool
    ) -> int:
        """Build permissions flag for PDF encryption."""
        # Start with no permissions
        perms = 0

        if allow_printing:
            perms |= 0b000000000100  # Bit 3: Print
        if allow_modification:
            perms |= 0b000000001000  # Bit 4: Modify contents
        if allow_copying:
            perms |= 0b000000010000  # Bit 5: Copy text and graphics

        return perms

    async def html_to_pdf(
        self,
        html_content: str = None,
        url: str = None,
    ) -> Tuple[bytes, int]:
        """
        Convert HTML content or URL to PDF.

        Args:
            html_content: HTML string to convert
            url: URL to fetch and convert (alternative to html_content)

        Returns:
            Tuple of (PDF bytes, estimated pages)
        """
        try:
            from weasyprint import HTML
            import asyncio

            if not html_content and not url:
                raise BadRequestError(message="Either html_content or url must be provided")

            # Run in thread pool (I/O-bound)
            loop = asyncio.get_event_loop()
            pdf_bytes = await loop.run_in_executor(
                _executor,
                lambda: self._convert_html_to_pdf(html_content, url)
            )

            # Estimate page count (rough estimate based on size)
            estimated_pages = max(1, len(pdf_bytes) // 50000)

            return pdf_bytes, estimated_pages

        except Exception as e:
            raise FileProcessingError(message=f"HTML to PDF conversion failed: {str(e)}")

    def _convert_html_to_pdf(self, html_content: str = None, url: str = None) -> bytes:
        """Synchronous HTML to PDF conversion."""
        from weasyprint import HTML

        if url:
            html = HTML(url=url)
        else:
            html = HTML(string=html_content)

        return html.write_pdf()

    async def from_docx(self, docx_bytes: bytes) -> Tuple[bytes, int]:
        """
        Convert DOCX to PDF using LibreOffice.

        Args:
            docx_bytes: DOCX file bytes

        Returns:
            Tuple of (PDF bytes, estimated pages)
        """
        import subprocess
        import shutil
        import asyncio

        # Check if LibreOffice is available
        if not shutil.which("libreoffice"):
            raise FileProcessingError(
                message="LibreOffice is not installed. Word to PDF conversion requires LibreOffice."
            )

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as docx_file:
            docx_file.write(docx_bytes)
            docx_path = docx_file.name

        pdf_path = docx_path.replace(".docx", ".pdf")

        try:
            # Run LibreOffice conversion in thread pool
            loop = asyncio.get_event_loop()
            success = await loop.run_in_executor(
                _executor,
                lambda: self._convert_docx_with_libreoffice(docx_path, pdf_path)
            )

            if not success or not os.path.exists(pdf_path):
                raise FileProcessingError(
                    message="DOCX conversion failed. The file may be corrupted or use unsupported features."
                )

            # Read result
            with open(pdf_path, "rb") as f:
                pdf_bytes = f.read()

            # Get page count
            reader = PdfReader(io.BytesIO(pdf_bytes))
            total_pages = len(reader.pages)

            return pdf_bytes, total_pages

        finally:
            # Cleanup
            if os.path.exists(docx_path):
                os.unlink(docx_path)
            if os.path.exists(pdf_path):
                os.unlink(pdf_path)

    def _convert_docx_with_libreoffice(self, docx_path: str, pdf_path: str) -> bool:
        """Convert DOCX to PDF using LibreOffice headless."""
        import subprocess

        try:
            # Get directory containing the DOCX file
            input_dir = os.path.dirname(docx_path)

            # Run LibreOffice in headless mode
            result = subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    input_dir,
                    docx_path,
                ],
                capture_output=True,
                timeout=60,
            )

            return result.returncode == 0

        except Exception:
            return False

    async def add_watermark_text(
        self,
        pdf_bytes: bytes,
        text: str,
        position: str = "center",
        opacity: float = 0.3,
        font_size: int = 40,
        rotation: int = 45,
        color: str = "#808080",
    ) -> Tuple[bytes, int]:
        """
        Add text watermark to PDF.

        Args:
            pdf_bytes: PDF file bytes
            text: Watermark text
            position: "center", "top", "bottom"
            opacity: 0.0 to 1.0
            font_size: Font size in points
            rotation: Rotation angle in degrees
            color: Hex color code

        Returns:
            Tuple of (watermarked PDF bytes, total pages)
        """
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.colors import HexColor

            # Parse color
            try:
                watermark_color = HexColor(color)
            except:
                watermark_color = HexColor("#808080")

            # Create watermark PDF
            watermark_buffer = io.BytesIO()
            c = canvas.Canvas(watermark_buffer, pagesize=letter)

            # Set watermark style with opacity
            color_with_alpha = HexColor(color, alpha=opacity)
            c.setFillColor(color_with_alpha)
            c.setFont("Helvetica-Bold", font_size)

            # Position watermark
            page_width, page_height = letter

            if position == "center":
                x, y = page_width / 2, page_height / 2
            elif position == "top":
                x, y = page_width / 2, page_height - 100
            elif position == "bottom":
                x, y = page_width / 2, 100
            else:
                x, y = page_width / 2, page_height / 2

            # Draw rotated text
            c.saveState()
            c.translate(x, y)
            c.rotate(rotation)
            c.drawCentredString(0, 0, text)
            c.restoreState()

            c.save()

            # Merge watermark with original PDF
            watermark_buffer.seek(0)
            watermark_reader = PdfReader(watermark_buffer)
            watermark_page = watermark_reader.pages[0]

            reader = PdfReader(io.BytesIO(pdf_bytes))
            writer = PdfWriter()

            for page in reader.pages:
                page.merge_page(watermark_page)
                writer.add_page(page)

            # Compress output
            writer.compress_identical_objects(remove_identicals=True, remove_orphans=True)

            output = io.BytesIO()
            writer.write(output)

            return output.getvalue(), len(reader.pages)

        except Exception as e:
            raise FileProcessingError(message=f"Watermark failed: {str(e)}")

    async def add_page_numbers(
        self,
        pdf_bytes: bytes,
        position: str = "bottom-center",
        format_string: str = "{page}",
        font_size: int = 10,
        start_page: int = 1,
        margin: int = 20,
    ) -> Tuple[bytes, int]:
        """
        Add page numbers to PDF.

        Args:
            pdf_bytes: PDF file bytes
            position: "top-left", "top-center", "top-right", "bottom-left", "bottom-center", "bottom-right"
            format_string: Format like "{page}", "Page {page}", "{page} of {total}"
            font_size: Font size in points
            start_page: Page number to start from
            margin: Margin from edge in points

        Returns:
            Tuple of (PDF with page numbers, total pages)
        """
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            total_pages = len(doc)

            # Parse position
            position_map = {
                "top-left": (margin, margin),
                "top-center": ("center", margin),
                "top-right": (-margin, margin),
                "bottom-left": (margin, -margin),
                "bottom-center": ("center", -margin),
                "bottom-right": (-margin, -margin),
            }

            base_x, base_y = position_map.get(position, ("center", -margin))

            for page_num in range(total_pages):
                page = doc[page_num]
                page_rect = page.rect

                # Format page number text
                text = format_string.format(
                    page=page_num + start_page,
                    total=total_pages
                )

                # Calculate position
                if base_x == "center":
                    x = page_rect.width / 2
                    align = fitz.TEXT_ALIGN_CENTER
                elif base_x < 0:
                    x = page_rect.width + base_x
                    align = fitz.TEXT_ALIGN_RIGHT
                else:
                    x = base_x
                    align = fitz.TEXT_ALIGN_LEFT

                if base_y < 0:
                    y = page_rect.height + base_y
                else:
                    y = base_y + font_size

                # Insert text
                page.insert_text(
                    fitz.Point(x, y),
                    text,
                    fontsize=font_size,
                    color=(0, 0, 0),
                    fontname="helv",
                    render_mode=0,
                )

            # Save with compression
            output = io.BytesIO()
            doc.save(output, garbage=4, deflate=True, clean=True)
            doc.close()

            return output.getvalue(), total_pages

        except Exception as e:
            raise FileProcessingError(message=f"Adding page numbers failed: {str(e)}")

    async def organize(
        self,
        pdf_bytes: bytes,
        page_order: list[int],
    ) -> Tuple[bytes, int]:
        """
        Organize PDF by reordering or deleting pages.

        Args:
            pdf_bytes: PDF file bytes
            page_order: List of page numbers (1-indexed) in desired order
                       e.g., [1, 3, 2] to swap pages 2 and 3
                       e.g., [1, 3] to remove page 2

        Returns:
            Tuple of (organized PDF bytes, total pages)
        """
        try:
            reader = PdfReader(io.BytesIO(pdf_bytes))
            total_pages = len(reader.pages)
            writer = PdfWriter()

            # Validate page numbers
            for page_num in page_order:
                if page_num < 1 or page_num > total_pages:
                    raise BadRequestError(
                        message=f"Invalid page number: {page_num}. PDF has {total_pages} pages."
                    )

            # Add pages in specified order
            for page_num in page_order:
                writer.add_page(reader.pages[page_num - 1])  # Convert to 0-indexed

            # Compress output
            writer.compress_identical_objects(remove_identicals=True, remove_orphans=True)

            output = io.BytesIO()
            writer.write(output)

            return output.getvalue(), len(page_order)

        except BadRequestError:
            raise
        except Exception as e:
            raise FileProcessingError(message=f"PDF organization failed: {str(e)}")

    async def crop(
        self,
        pdf_bytes: bytes,
        left: int = 0,
        top: int = 0,
        right: int = 0,
        bottom: int = 0,
        pages: str = "all",
    ) -> Tuple[bytes, int]:
        """
        Crop PDF pages by removing margins.

        Args:
            pdf_bytes: PDF file bytes
            left: Left margin to crop in points
            top: Top margin to crop in points
            right: Right margin to crop in points
            bottom: Bottom margin to crop in points
            pages: Pages to crop - "all", "even", "odd", or "1,3,5-7"

        Returns:
            Tuple of (cropped PDF bytes, total pages)
        """
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            total_pages = len(doc)

            # Parse page selection
            if pages.lower() == "all":
                pages_to_crop = set(range(total_pages))
            elif pages.lower() == "even":
                pages_to_crop = set(range(1, total_pages, 2))
            elif pages.lower() == "odd":
                pages_to_crop = set(range(0, total_pages, 2))
            else:
                pages_to_crop = self._parse_page_selection(pages, total_pages)

            # Crop pages
            for page_num in range(total_pages):
                if page_num in pages_to_crop:
                    page = doc[page_num]
                    rect = page.rect

                    # Create new crop box
                    crop_rect = fitz.Rect(
                        rect.x0 + left,
                        rect.y0 + top,
                        rect.x1 - right,
                        rect.y1 - bottom,
                    )

                    # Validate crop box
                    if crop_rect.width <= 0 or crop_rect.height <= 0:
                        raise BadRequestError(
                            message=f"Invalid crop dimensions for page {page_num + 1}. Crop box is too large."
                        )

                    page.set_cropbox(crop_rect)

            # Save with compression
            output = io.BytesIO()
            doc.save(output, garbage=4, deflate=True, clean=True)
            doc.close()

            return output.getvalue(), total_pages

        except BadRequestError:
            raise
        except Exception as e:
            raise FileProcessingError(message=f"PDF crop failed: {str(e)}")

    async def from_excel(self, excel_bytes: bytes) -> Tuple[bytes, int]:
        """
        Convert Excel (XLSX) to PDF using LibreOffice.

        Args:
            excel_bytes: Excel file bytes

        Returns:
            Tuple of (PDF bytes, estimated pages)
        """
        import subprocess
        import shutil
        import asyncio

        # Check if LibreOffice is available
        if not shutil.which("libreoffice"):
            raise FileProcessingError(
                message="LibreOffice is not installed. Excel to PDF conversion requires LibreOffice."
            )

        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as excel_file:
            excel_file.write(excel_bytes)
            excel_path = excel_file.name

        pdf_path = excel_path.replace(".xlsx", ".pdf")

        try:
            # Run LibreOffice conversion in thread pool
            loop = asyncio.get_event_loop()
            success = await loop.run_in_executor(
                _executor, lambda: self._convert_office_with_libreoffice(excel_path, pdf_path)
            )

            if not success or not os.path.exists(pdf_path):
                raise FileProcessingError(
                    message="Excel conversion failed. The file may be corrupted or use unsupported features."
                )

            # Read result
            with open(pdf_path, "rb") as f:
                pdf_bytes = f.read()

            # Get page count
            reader = PdfReader(io.BytesIO(pdf_bytes))
            total_pages = len(reader.pages)

            return pdf_bytes, total_pages

        finally:
            # Cleanup
            if os.path.exists(excel_path):
                os.unlink(excel_path)
            if os.path.exists(pdf_path):
                os.unlink(pdf_path)

    async def from_powerpoint(self, pptx_bytes: bytes) -> Tuple[bytes, int]:
        """
        Convert PowerPoint (PPTX) to PDF using LibreOffice.

        Args:
            pptx_bytes: PowerPoint file bytes

        Returns:
            Tuple of (PDF bytes, estimated pages)
        """
        import subprocess
        import shutil
        import asyncio

        # Check if LibreOffice is available
        if not shutil.which("libreoffice"):
            raise FileProcessingError(
                message="LibreOffice is not installed. PowerPoint to PDF conversion requires LibreOffice."
            )

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as pptx_file:
            pptx_file.write(pptx_bytes)
            pptx_path = pptx_file.name

        pdf_path = pptx_path.replace(".pptx", ".pdf")

        try:
            # Run LibreOffice conversion in thread pool
            loop = asyncio.get_event_loop()
            success = await loop.run_in_executor(
                _executor, lambda: self._convert_office_with_libreoffice(pptx_path, pdf_path)
            )

            if not success or not os.path.exists(pdf_path):
                raise FileProcessingError(
                    message="PowerPoint conversion failed. The file may be corrupted or use unsupported features."
                )

            # Read result
            with open(pdf_path, "rb") as f:
                pdf_bytes = f.read()

            # Get page count
            reader = PdfReader(io.BytesIO(pdf_bytes))
            total_pages = len(reader.pages)

            return pdf_bytes, total_pages

        finally:
            # Cleanup
            if os.path.exists(pptx_path):
                os.unlink(pptx_path)
            if os.path.exists(pdf_path):
                os.unlink(pdf_path)

    def _convert_office_with_libreoffice(self, input_path: str, output_path: str) -> bool:
        """Convert Office files (DOCX, XLSX, PPTX) to PDF using LibreOffice headless."""
        import subprocess

        try:
            # Get directory containing the input file
            input_dir = os.path.dirname(input_path)

            # Run LibreOffice in headless mode
            result = subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    input_dir,
                    input_path,
                ],
                capture_output=True,
                timeout=60,
            )

            return result.returncode == 0

        except Exception:
            return False
    async def to_excel(self, pdf_bytes: bytes) -> Tuple[bytes, int]:
        """
        Convert PDF to Excel (XLSX) with Ultra-Fidelity (Gold Standard).
        Reconstructs page layouts, embeds images at exact coordinates, and maps vector colors.
        """
        try:
            import pandas as pd
            import fitz
            import openpyxl
            from openpyxl.styles import PatternFill, Font, Alignment
            from openpyxl.utils import get_column_letter
            from openpyxl.drawing.image import Image as XLImage
            import tempfile
            import os
            
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            total_pages = len(doc)
            
            output = io.BytesIO()
            workbook = openpyxl.Workbook()
            if "Sheet" in workbook.sheetnames:
                workbook.remove(workbook["Sheet"])
            
            # Helper to convert fitz color to hex
            def fitz_to_hex(color):
                if color is None: return None
                if isinstance(color, (list, tuple)):
                    rgb = color
                else:
                    # sRGB integer
                    rgb = [(color >> 16) & 0xFF, (color >> 8) & 0xFF, color & 0xFF]
                    return '{:02x}{:02x}{:02x}'.format(rgb[0], rgb[1], rgb[2]).upper()
                
                return '{:02x}{:02x}{:02x}'.format(
                    int(rgb[0] * 255 if max(rgb) <= 1.0 else rgb[0]), 
                    int(rgb[1] * 255 if max(rgb) <= 1.0 else rgb[1]), 
                    int(rgb[2] * 255 if max(rgb) <= 1.0 else rgb[2])
                ).upper()

            for page_num, page in enumerate(doc):
                # --- SHEET 1: RECONSTRUCTED PAGE LAYOUT (Ultra-Fidelity) ---
                main_ws = workbook.create_sheet(title=f"Page {page_num+1}")
                
                # Set default column and row sizes for a tight grid
                COL_WIDTH = 10 # chars
                ROW_HEIGHT = 15 # points
                
                # 1pt ~ 1/72 inch. Excel units are different. 
                # Approx mapping: 1pt in PDF -> 1pt row height in Excel.
                # 1pt in PDF -> approx 0.15 character width in Excel.
                
                drawings = page.get_drawings()
                blocks = page.get_text("dict")["blocks"]
                
                # --- A. BACKGROUND SHAPES & COLORS ---
                # We map drawings (rects/fills) to the grid
                for draw in drawings:
                    if "fill" in draw and draw["fill"]:
                        rect = draw["rect"]
                        bg_color = fitz_to_hex(draw["fill"])
                        
                        # Map rect to cell range
                        s_row = int(rect.y0 / ROW_HEIGHT) + 1
                        e_row = int(rect.y1 / ROW_HEIGHT) + 1
                        s_col = int(rect.x0 / (COL_WIDTH * 5)) + 1 # approx scale factor
                        e_col = int(rect.x1 / (COL_WIDTH * 5)) + 1
                        
                        for r in range(s_row, e_row + 1):
                            for c in range(s_col, e_col + 1):
                                cell = main_ws.cell(row=r, column=c)
                                cell.fill = PatternFill(start_color=bg_color, end_color=bg_color, fill_type="solid")

                # --- B. IMAGES ---
                image_list = page.get_images(full=True)
                for img_idx, img_info in enumerate(image_list):
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # Get image position on page
                    # Note: get_image_info (v1.18.0+) or page.get_image_rects
                    rects = page.get_image_rects(xref)
                    if rects:
                        rect = rects[0]
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                            tmp.write(image_bytes)
                            tmp_path = tmp.name
                        
                        try:
                            img = XLImage(tmp_path)
                            # Approximate position
                            # anchor: 'A1' etc.
                            c_idx = int(rect.x0 / (COL_WIDTH * 5)) + 1
                            r_idx = int(rect.y0 / ROW_HEIGHT) + 1
                            
                            # Scale image to fit the bbox
                            # Excel Image width/height are in pixels. PDF pts are 1/72.
                            # 1 pt = 1.333 pixels
                            img.width = rect.width * 1.333
                            img.height = rect.height * 1.333
                            
                            main_ws.add_image(img, get_column_letter(c_idx) + str(r_idx))
                        finally:
                            if os.path.exists(tmp_path):
                                os.remove(tmp_path)

                # --- C. TEXT OVERLAYS ---
                # Re-map text to cells atop icons/colors
                for block in blocks:
                    if block["type"] == 0: # Text
                        for line in block["lines"]:
                            for span in line["spans"]:
                                bbox = span["bbox"]
                                r_idx = int(bbox[1] / ROW_HEIGHT) + 1
                                c_idx = int(bbox[0] / (COL_WIDTH * 5)) + 1
                                
                                cell = main_ws.cell(row=r_idx, column=c_idx, value=span["text"])
                                
                                # Font Styles
                                is_bold = bool(span["flags"] & 16)
                                is_italic = bool(span["flags"] & 2)
                                hex_color = fitz_to_hex(span["color"])
                                cell.font = Font(
                                    bold=is_bold, 
                                    italic=is_italic, 
                                    color=hex_color if hex_color != "000000" else None, 
                                    size=span["size"]
                                )
                                cell.alignment = Alignment(vertical="center")

                # --- SHEET 2: STRUCTURED TABLES (Analysis Layer) ---
                tabs = page.find_tables()
                for i, tab in enumerate(tabs.tables):
                    sheet_name = f"P{page_num+1} Table {i+1}"
                    ws = workbook.create_sheet(title=sheet_name[:31])
                    
                    data = tab.extract()
                    for r_idx, row in enumerate(data):
                        for c_idx, value in enumerate(row):
                            cell = ws.cell(row=r_idx + 1, column=c_idx + 1, value=value)
                            
                            # Extract background color for this specific table cell
                            cell_bbox = tab.cells[r_idx * tab.col_count + c_idx]
                            if cell_bbox:
                                for draw in drawings:
                                    if "fill" in draw and draw["fill"]:
                                        d_rect = draw["rect"]
                                        if d_rect.intersects(cell_bbox) and d_rect.width > 2 and d_rect.height > 2:
                                            bg_color = fitz_to_hex(draw["fill"])
                                            cell.fill = PatternFill(start_color=bg_color, end_color=bg_color, fill_type="solid")

                    # Column widths for the data sheet
                    for column_cells in ws.columns:
                        length = max(len(str(cell.value or "")) for cell in column_cells)
                        ws.column_dimensions[get_column_letter(column_cells[0].column)].width = min(length + 2, 60)

                # Set grid dimensions for main sheet
                for i in range(1, 100):
                    main_ws.column_dimensions[get_column_letter(i)].width = COL_WIDTH / 2 # tighter grid
                for i in range(1, 200):
                    main_ws.row_dimensions[i].height = ROW_HEIGHT

            # Finalize metadata
            props = workbook.properties
            props.title = "Ultra-Fidelity Conversion"
            props.creator = "Tulz"
            props.description = "Premium PDF to Excel conversion with images, styles, and vector fidelity."
            
            doc.close()
            workbook.save(output)
            return output.getvalue(), total_pages

        except Exception as e:
            import traceback
            print(f"DEBUG: Excel Failure: {traceback.format_exc()}")
            raise FileProcessingError(message=f"PDF to Excel Ultra-Fidelity conversion failed: {str(e)}")

    async def to_powerpoint(self, pdf_bytes: bytes) -> Tuple[bytes, int]:
        """
        Convert PDF to PowerPoint (PPTX) with iLovePDF-style high fidelity.
        Renders a text-free background to preserve design/graphics and 
        overlays editable text boxes at exact coordinates.
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            import fitz
            import io
            
            # Load original doc
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            total_pages = len(doc)
            
            prs = Presentation()
            
            # Set document properties (metadata) to help OS identify the file type
            props = prs.core_properties
            props.title = "Converted Presentation"
            props.author = "Tulz"
            props.category = "Conversion"
            props.comments = "Generated by Tulz PDF to PowerPoint Converter"

            blank_slide_layout = prs.slide_layouts[6]

            # Process each page
            for page_num in range(total_pages):
                page = doc[page_num]
                
                # 1. Create a "clean" version of the page for background (no text)
                # We do this by creating a copy and redacting all text
                temp_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
                temp_page = temp_doc[page_num]
                
                # Find all text and redact it (render as clean background)
                for block in temp_page.get_text("dict")["blocks"]:
                    if block["type"] == 0:  # Text block
                        for line in block["lines"]:
                            for span in line["spans"]:
                                temp_page.add_redact_annot(span["bbox"])
                
                temp_page.apply_redactions(graphics=0) # Only remove text
                
                # Render the clean page to high-res image
                pix = temp_page.get_pixmap(matrix=fitz.Matrix(300/72, 300/72)) # 300 DPI
                img_data = pix.tobytes("png")
                temp_doc.close()
                
                # 2. Reconstruct slide
                slide_width = Inches(page.rect.width / 72.0)
                slide_height = Inches(page.rect.height / 72.0)
                
                # Match slide size to PDF page
                prs.slide_width = slide_width
                prs.slide_height = slide_height
                
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Set the clean background image
                img_stream = io.BytesIO(img_data)
                slide.shapes.add_picture(img_stream, 0, 0, width=slide_width, height=slide_height)
                
                # 3. Overlay editable text boxes
                for block in page.get_text("dict")["blocks"]:
                    if block["type"] == 0:
                        for line in block["lines"]:
                            for span in line["spans"]:
                                text = span["text"]
                                if not text.strip():
                                    continue
                                
                                bbox = span["bbox"]
                                left = Inches(bbox[0] / 72.0)
                                top = Inches(bbox[1] / 72.0)
                                width = Inches((bbox[2] - bbox[0]) / 72.0)
                                height = Inches((bbox[3] - bbox[1]) / 72.0)
                                
                                tx_box = slide.shapes.add_textbox(left, top, width, height)
                                tf = tx_box.text_frame
                                tf.word_wrap = False # Keep layout tight
                                
                                p = tf.paragraphs[0]
                                p.text = text
                                p.font.size = Pt(span["size"])
                                p.font.bold = bool(span["flags"] & 16)
                                p.font.italic = bool(span["flags"] & 2)
                                
                                # Match color if possible
                                try:
                                    color = span["color"] # sRGB integer
                                    r = (color >> 16) & 0xFF
                                    g = (color >> 8) & 0xFF
                                    b = color & 0xFF
                                    p.font.color.rgb = RGBColor(r, g, b)
                                except:
                                    pass

            doc.close()

            output = io.BytesIO()
            prs.save(output)
            
            return output.getvalue(), total_pages

        except Exception as e:
            import traceback
            print(f"DEBUG: PPTX Failure: {traceback.format_exc()}")
            raise FileProcessingError(message=f"High-fidelity PDF to PowerPoint conversion failed: {str(e)}")
